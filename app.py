import streamlit as st
from PyPDF2 import PdfReader
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain.text_splitter import CharacterTextSplitter
from langchain.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
import google.generativeai as genai  # Import the google.generativeai module
import os
import pandas as pd
from pptx import Presentation
from docx import Document
import time as t
from dotenv import load_dotenv

load_dotenv()

GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=GOOGLE_API_KEY)

MAX_FILE_SIZE_MB = 2
DOCUMENTS_FOLDER = "admin_handle_files"  # Change this to match your folder name

# SessionState class to handle persistent states
class SessionState:
    def __init__(self, **kwargs):
        self.__dict__.update(kwargs)

# Process documents to extract and chunk text for embedding
def process_documents(documents_folder):
    documents = []
    for file_name in os.listdir(documents_folder):
        file_path = os.path.join(documents_folder, file_name)
        if file_name.endswith(".pdf"):
            documents.extend([page.extract_text() for page in PdfReader(file_path).pages])
        elif file_name.endswith(".txt"):
            with open(file_path, "r") as file:
                documents.append(file.read())
        elif file_name.endswith(".docx"):
            doc = Document(file_path)
            documents.extend([para.text for para in doc.paragraphs])
        elif file_name.endswith(".xlsx"):
            for sheet in pd.ExcelFile(file_path).sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet)
                documents.extend(df.astype(str).stack().dropna().tolist())
        elif file_name.endswith(".pptx"):
            presentation = Presentation(file_path)
            documents.extend([shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")])

    text = "\n".join(documents)
    if not text:
        return None

    text_chunks = CharacterTextSplitter(separator="\n", chunk_size=10000, chunk_overlap=500).split_text(text)
    if not text_chunks:
        return None

    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    return FAISS.from_texts(text_chunks, embeddings)

# Function to handle query and response generation
def answer_query(query, docsearch, chain):
    docs = docsearch.similarity_search(query)
    return chain.run(input_documents=docs, question=query)

# Load the QA Chain with the specified model and prompt template
def get_chain():
    prompt_template = """
        Provide a thorough answer based on the given context. 
        If the required information is not found within the context, 
        respond with 'The answer is not available in the context.'
        Context:\n{context}\n
        Question:\n{question}\n
        Answer:
    """
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    return load_qa_chain(model, chain_type="stuff", prompt=prompt)

# Admin functionality for document upload and management
def admin_panel(session_state):
    if st.sidebar.button("Upload Documents"):
        uploaded_files = st.file_uploader("Upload your documents below", accept_multiple_files=True, type=["pdf", "txt", "docx", "xlsx", "pptx"])
        if uploaded_files:
            if not os.path.exists(DOCUMENTS_FOLDER):
                os.makedirs(DOCUMENTS_FOLDER)

            for file in uploaded_files:
                if len(file.getvalue()) > MAX_FILE_SIZE_MB * 1024 * 1024:
                    st.error(f"File size exceeds {MAX_FILE_SIZE_MB} MB limit: {file.name}")
                    continue
                with open(os.path.join(DOCUMENTS_FOLDER, file.name), "wb") as f:
                    f.write(file.getvalue())
            session_state.docsearch = process_documents(DOCUMENTS_FOLDER)

    if session_state.docsearch:
        display_chat_interface(session_state)

    display_uploaded_documents()

# Normal user interface
def user_panel(session_state):
    if os.path.exists(DOCUMENTS_FOLDER):
        session_state.docsearch = process_documents(DOCUMENTS_FOLDER)
        if session_state.docsearch:
            display_chat_interface(session_state)
        display_uploaded_documents()
    else:
        st.warning("No documents are currently present. Please wait for the admin to initiate upload")

# Display uploaded documents with an option to delete (for admins)
def display_uploaded_documents():
    if os.path.exists(DOCUMENTS_FOLDER):
        uploaded_documents = os.listdir(DOCUMENTS_FOLDER)
        if uploaded_documents:
            st.sidebar.header("Uploaded Documents:")
            selected_documents = [file for file in uploaded_documents if st.sidebar.checkbox(file)]
            if st.sidebar.button("Delete Selected Documents"):
                for file_name in selected_documents:
                    os.remove(os.path.join(DOCUMENTS_FOLDER, file_name))
                st.sidebar.success("Documents deleted successfully. Refresh to update the list.")

# Display chat interface
def display_chat_interface(session_state):
    st.subheader("AI bot")
    user_input = st.text_input("You:", key="user_input")
    if st.button("Send"):
        response = answer_query(user_input, session_state.docsearch, session_state.chain)
        st.text("Bot:")
        with st.spinner("Wait for a moment..."):
            t.sleep(2)
            st.write(response)

# Main function for Streamlit UI
def main():
    st.set_page_config(page_title="Document Search Bot")
    
    # Custom CSS and title
    st.markdown("""
        <style>
        @import url('https://fonts.googleapis.com/css2?family=Londrina+Sketch&display=swap');
        @import url('https://fonts.googleapis.com/css2?family=Quicksand&display=swap');

        /* Global background and text color */
        body {
            background-color: black;
            color: white;
        }

        .custom-title {
            font-family: 'Londrina Sketch', cursive;
            font-size: 48px;
            color: black;
            text-align: center;
        }

        /* Sidebar styles */
        .css-1d391kg {
            background-color: #01083b;
        }

        /* Sidebar text color */
        .css-1d391kg, .css-1d391kg .stSidebar .stSidebar-content {
            font-family: 'Quicksand', sans-serif;
            color: white !important;
        }

        .stSidebar [data-testid="stSidebar"] {
            background-color: #01083b;
        }

        /* Uploaded documents text color */
        .sidebar .sidebar-header {
            color: gold !important;
        }
        </style>
        <h1 class="custom-title">DOCUMENT SEARCH BOT</h1>
    """, unsafe_allow_html=True)

    # Ensure the documents folder exists
    if not os.path.exists(DOCUMENTS_FOLDER):
        os.makedirs(DOCUMENTS_FOLDER)

    session_state = SessionState(docsearch=None, chain=get_chain(), is_admin=False)

    # Move "Select user type" to the main window
    user_type = st.radio("Select user type:", ["User", "Admin"], index=0)

    if user_type == "Admin":
        admin_password = st.text_input("Enter admin password:", type="password")
        if admin_password == "admin@123":
            session_state.is_admin = True
            st.success("Authentication successful")
        else:
            st.error("Incorrect password.")

    if session_state.is_admin:
        admin_panel(session_state)
    else:
        user_panel(session_state)

    if st.sidebar.button("Refresh Page"):
        st.experimental_rerun()

if __name__ == '__main__':
    main()
